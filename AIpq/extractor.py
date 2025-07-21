import  uuid
from datetime import datetime
from fastapi import APIRouter, File, UploadFile, HTTPException
from fastapi.responses import JSONResponse
from pptx import Presentation
from docx import Document
import pdfplumber
from pdf2image import convert_from_path
import subprocess
import tempfile
import os
import shutil
import io
import whisper
import cv2
import numpy as np
from PIL import Image
import easyocr

from aiapiclient import client

router = APIRouter()

# ---------- 统一 OCR ----------
ocr = easyocr.Reader(['ch_sim', 'en'], gpu=False)

# ---------- 工具函数 ----------
def extract_txt(p):
    return open(p, encoding="utf-8").read()

def extract_pptx(p):
    from pptx.enum.shapes import MSO_SHAPE_TYPE
    prs = Presentation(p)
    lines = []
    for s in prs.slides:
        lines += [sh.text.strip() for sh in s.shapes if hasattr(sh, "text") and sh.text.strip()]
        if s.has_notes_slide and s.notes_slide.notes_text_frame:
            lines.append(s.notes_slide.notes_text_frame.text.strip())
        for sh in s.shapes:
            if sh.shape_type == MSO_SHAPE_TYPE.PICTURE:
                img = Image.open(io.BytesIO(sh.image.blob)).convert("RGB")
                res = ocr.readtext(np.array(img), detail=0)
                lines.extend(res)
    return "\n".join(lines)

def extract_ppt(p):


    tmp = tempfile.mkdtemp()
    try:
        subprocess.run([
            r"D:\libre\program\soffice.exe",  # 绝对路径
            "--headless",
            "--convert-to", "pptx",
            "--outdir", tmp,
            p
        ], check=True, stdout=subprocess.DEVNULL, stderr=subprocess.DEVNULL)

        pptx = os.path.join(tmp, os.path.splitext(os.path.basename(p))[0] + ".pptx")
        text = extract_pptx(pptx) if os.path.exists(pptx) else ""
    finally:
        shutil.rmtree(tmp, ignore_errors=True)

    return text

def extract_pdf(p):
    full_text = []
    with pdfplumber.open(p) as pdf:
        for page in pdf.pages:
            txt = page.extract_text() or ""
            if txt.strip():
                full_text.append(txt)
            else:
                img = convert_from_path(
                    p, dpi=300, first_page=page.page_number, last_page=page.page_number, thread_count=1
                )[0]
                res = ocr.readtext(np.array(img), detail=0)
                full_text.extend(res)
    return "\n".join(full_text)

def extract_docx(p):

    doc = Document(p)
    lines = [para.text.strip() for para in doc.paragraphs if para.text.strip()]

    # 表格文字
    for tbl in doc.tables:
        for row in tbl.rows:
            for cell in row.cells:
                txt = cell.text.strip()
                if txt:
                    lines.append(txt)

    # 遍历所有段落中的图片
    for rel in doc.part.rels.values():
        if "image" in rel.target_ref:
            try:
                img_blob = rel.target_part.blob
                img = Image.open(io.BytesIO(img_blob)).convert("RGB")
                ocr_lines = ocr.readtext(np.array(img), detail=0)
                lines.extend(ocr_lines)
            except Exception:
                continue   # 图片解析失败直接跳过

    return "\n".join(lines)

def extract_doc(p):
    import subprocess, tempfile, os, shutil
    tmp = tempfile.mkdtemp()
    subprocess.run([
        r"D:\libre\program\soffice.exe",  # 你本地的路径
        "--headless",
        "--convert-to", "docx",
        "--outdir", tmp,
        p
    ], check=True, stdout=subprocess.DEVNULL, stderr=subprocess.DEVNULL)

    docx = os.path.join(tmp, os.path.splitext(os.path.basename(p))[0] + ".docx")
    text = extract_docx(docx) if os.path.exists(docx) else ""
    shutil.rmtree(tmp, ignore_errors=True)
    return text

def extract_audio(p: str) -> str:
    """
    使用 openai-whisper 本地转写音频
    支持 mp3/wav/m4a/flac/ogg 等常见格式
    """
    model = whisper.load_model("base")  # 可选 tiny/base/small/medium/large
    result = model.transcribe(p, language="zh")  # language="zh" 可提升中文准确率
    return result["text"].strip()



def extract_video(
    video_path: str,
    skip_fps: int = 2,     # 每秒抽几帧，调大可提速
    dedup: bool = True     # 是否去重
) -> str:
    """
    把视频逐帧 OCR，返回所有出现的文字（去重后）
    """
    reader = easyocr.Reader( ["ch_sim","en"], gpu=False)   # GPU=True 更快
    cap = cv2.VideoCapture(video_path)
    fps = cap.get(cv2.CAP_PROP_FPS) or 25
    frame_id = 0
    seen = set()
    all_text = []

    while True:
        ret, frame = cap.read()
        if not ret:
            break

        # 按 skip_fps 抽帧
        if frame_id % max(1, int(fps // skip_fps)) != 0:
            frame_id += 1
            continue

        # 转 PIL，再转灰度
        img = Image.fromarray(cv2.cvtColor(frame, cv2.COLOR_BGR2RGB)).convert("L")

        # OCR
        texts = reader.readtext(np.array(img), detail=0)
        for t in texts:
            t = t.strip()
            if t:
                if dedup and t in seen:
                    continue
                seen.add(t)
                all_text.append(t)

        frame_id += 1

    cap.release()
    return polish_text("\n".join(all_text))

#整理文本
def polish_text(raw: str) -> str:
    if not raw.strip():
        return ""

    prompt = (
        "你是一位中文文本整理专家，请对下列 OCR/转写结果进行清洗：\n"
        "- 删除乱码、特殊符号、无意义字符；\n"
        "- 合并重复句子，保留一份；\n"
        "- 调整语序、断句，使语义通顺；\n"
        "- 保留原文中的换行符（\\n），不要把多行合并成一段；\n"
        "- 不增删原文信息，仅做整理；\n"
        "- 直接输出整理后的纯文本，不解释。\n\n"
        f"待整理文本：\n{raw[:4000]}"
    )

    resp = client.chat.completions.create(
        model="gpt-4.1-mini",
        messages=[{"role": "user", "content": prompt}],
        max_tokens=800,
        temperature=0.3,
        timeout=15
    )
    return resp.choices[0].message.content.strip()

# ---------- FastAPI 路由 ----------


@router.post("/extract")
async def extract_text(file: UploadFile = File(...)):
    # 1. 校验扩展名
    ext = os.path.splitext(file.filename)[-1].lower()
    if ext not in {
        ".txt", ".ppt", ".pptx", ".pdf",
        ".doc", ".docx",
        ".wav", ".mp3", ".m4a", ".flac", ".ogg",".mp4",".avi",".mov",".mkv"
    }:
        raise HTTPException(status_code=400, detail="Unsupported file type")

    # 2. 保存上传文件到临时目录
    tmp = os.path.join(tempfile.gettempdir(), f"{uuid.uuid4()}{ext}")
    with open(tmp, "wb") as buffer:
        shutil.copyfileobj(file.file, buffer)
    try:
        # 3. 根据扩展名调用对应解析函数
        text = {
            ".txt":  extract_txt,
            ".pptx": extract_pptx,
            ".ppt":  extract_ppt,
            ".pdf":  extract_pdf,
            ".docx": extract_docx,
            ".doc":  extract_doc,
            ".wav":  extract_audio,
            ".mp3":  extract_audio,
            ".m4a":  extract_audio,
            ".flac": extract_audio,
            ".ogg":  extract_audio,
            ".mp4": extract_video,
            ".avi": extract_video,
            ".mov": extract_video,
            ".mkv": extract_video,
        }[ext](tmp)

    finally:
        # 4. 清理临时文件
        os.remove(tmp)

    # 5. 返回结果
    return JSONResponse(
        content={
            "filename": file.filename,
            "text": text,
            "length": len(text),
            "timestamp": datetime.now().isoformat(),
        }
    )