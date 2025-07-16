# text_extractor_api.py
import os, uuid, re, shutil, tempfile, threading
from datetime import datetime
from fastapi import FastAPI, File, UploadFile, HTTPException
from fastapi.responses import JSONResponse




# -------------------------------------------------
# OCR 单例
# -------------------------------------------------
_ocr = None
_lock = threading.Lock()

def get_ocr():
    """线程安全的 PaddleOCR 单例"""
    global _ocr
    if _ocr is None:
        with _lock:
            if _ocr is None:
                from paddleocr import PaddleOCR
                _ocr = PaddleOCR(use_textline_orientation=True, lang='ch')
    return _ocr

# -------------------------------------------------
# 通用文本提取工具
# -------------------------------------------------
def extract_txt(file_path: str) -> str:
    with open(file_path, encoding="utf-8") as f:
        return f.read()

def extract_pptx(file_path: str) -> str:
    from pptx import Presentation
    from pptx.enum.shapes import MSO_SHAPE_TYPE
    ocr = get_ocr()

    prs = Presentation(file_path)
    lines = []
    for slide in prs.slides:
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                lines.append(shape.text)
        if slide.has_notes_slide and slide.notes_slide.notes_text_frame:
            notes = slide.notes_slide.notes_text_frame.text
            if notes.strip():
                lines.append(notes)
        for shape in slide.shapes:
            if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
                with tempfile.NamedTemporaryFile(suffix=".png", delete=False) as tmp:
                    tmp.write(shape.image.blob)
                    tmp.flush()
                    result = ocr.predict(tmp.name)
                    if result and result[0]:
                        for item in result[0]:
                            lines.append(item[1][0])
                os.remove(tmp.name)
    return "\n".join(lines)

def extract_ppt(file_path: str) -> str:
    import olefile
    from pdf2image import convert_from_bytes
    ocr = get_ocr()

    lines = []
    with olefile.OleFileIO(file_path) as ole:
        # 文本流
        for name in ole.listdir():
            if name[0].lower().startswith(('powerpoint document', 'notes')):
                data = ole.openstream(name).read()
                texts = re.findall(rb'(?:[\x20-\x7e])\x00(?:[\x20-\x7e])\x00', data)
                for t in texts:
                    try:
                        lines.append(t.decode('utf-16le'))
                    except UnicodeDecodeError:
                        pass
        # 图片 OCR
        for name in ole.listdir():
            if name[0] == 'Pictures':
                pic_data = ole.openstream(name).read()
                image = convert_from_bytes(pic_data, dpi=300)[0]
                result = ocr.predict(image)
                if result and result[0]:
                    for item in result[0]:
                        lines.append(item[1][0])
    return "\n".join(lines)

def extract_pdf(file_path: str) -> str:
    import pdfplumber
    text = ""
    with pdfplumber.open(file_path) as pdf:
        for page in pdf.pages:
            text += page.extract_text() or ""
    if text.strip():
        return text
    else:  # 扫描型
        from pdf2image import convert_from_path
        ocr = get_ocr()
        images = convert_from_path(file_path, dpi=300)
        lines = []
        for img in images:
            result = ocr.predict(img)
            if result and result[0]:
                for item in result[0]:
                    lines.append(item[1][0])
        return "\n".join(lines)

def extract_docx(file_path: str) -> str:
    from docx import Document
    ocr = get_ocr()

    doc = Document(file_path)
    lines = [p.text for p in doc.paragraphs if p.text.strip()]

    # 表格文字
    for table in doc.tables:
        for row in table.rows:
            for cell in row.cells:
                txt = cell.text.strip()
                if txt:
                    lines.append(txt)

    # 内嵌图片 OCR
    for rel in doc.part.rels.values():
        if "image" in rel.target_ref:
            img_data = rel.target_part.blob
            with tempfile.NamedTemporaryFile(suffix=".png", delete=False) as tmp:
                tmp.write(img_data)
                tmp.flush()
                result = ocr.predict(tmp.name)
                if result and result[0]:
                    for item in result[0]:
                        lines.append(item[1][0])
            os.remove(tmp.name)
    return "\n".join(lines)

def extract_doc(file_path: str) -> str:
    import olefile
    from pdf2image import convert_from_bytes
    ocr = get_ocr()

    lines = []
    with olefile.OleFileIO(file_path) as ole:
        # 文字流
        for name in ole.listdir():
            if name[0].lower().startswith(('worddocument', '1table')):
                data = ole.openstream(name).read()
                texts = re.findall(rb'[\x20-\x7e]{3,}', data)
                for t in texts:
                    try:
                        lines.append(t.decode('latin-1'))
                    except UnicodeDecodeError:
                        pass
        # 图片 OCR
        for name in ole.listdir():
            if name[0].lower().startswith('pictures'):
                img_data = ole.openstream(name).read()
                image = convert_from_bytes(img_data, dpi=300)[0]
                result = ocr.predict(image)
                if result and result[0]:
                    for item in result[0]:
                        lines.append(item[1][0])
    return "\n".join(lines)


# FastAPI
# -------------------------------------------------
app = FastAPI(title="Text Extractor API")


@app.post("/extract")
async def extract_text(file: UploadFile = File(...)):
    ext = os.path.splitext(file.filename)[-1].lower()
    if ext not in {".txt", ".ppt", ".pptx", ".pdf", ".doc", ".docx"}:
        raise HTTPException(status_code=400, detail="Unsupported file type")

    tmp_root = os.getenv("TMP_DIR", tempfile.gettempdir())
    os.makedirs(tmp_root, exist_ok=True)
    tmp_path = os.path.join(tmp_root, f"{uuid.uuid4()}{ext}")

    with open(tmp_path, "wb") as f:
        shutil.copyfileobj(file.file, f)

    try:
        if ext == ".txt":
            text = extract_txt(tmp_path)
        elif ext == ".ppt":
            text = extract_ppt(tmp_path)
        elif ext == ".pptx":
            text = extract_pptx(tmp_path)
        elif ext == ".pdf":
            text = extract_pdf(tmp_path)
        elif ext == ".doc":
            text = extract_doc(tmp_path)
        elif ext == ".docx":
            text = extract_docx(tmp_path)
    except Exception as e:
        raise HTTPException(status_code=500, detail=str(e))
    finally:
        os.remove(tmp_path)

    created_at = datetime.now()


    return JSONResponse(content={
        "filename": file.filename,
        "text": text,
        "length": len(text),
        "timestamp": created_at.isoformat()
    })

if __name__ == "__main__":
    import uvicorn
    uvicorn.run(app, host="0.0.0.0", port=8000)